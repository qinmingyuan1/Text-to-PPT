import streamlit as st
import base64
import pptx
from pptx.util import Inches, Pt
import os
import toml,json, requests

import pandas as pd
import docx2txt
import comtypes.client

# popup modal
from streamlit_modal import Modal

#** debug 
# import time
# import tempfile

#openai.api_key = os.getenv('OPENAI_API_KEY')  # Replace with your actual API key
file_path = './credential.txt'
if os.path.exists(file_path):
    with open(file_path, 'r') as f:
        secrets = toml.load(f)
else:
    st.warning("Credentials file not found. Please upload the credentials file.")
OPENROUTER_API_KEY = secrets['OPENROUTER']['OPENROUTER_API_KEY']
# Define custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

def generate_slide_titles(MinSlide,MaxSlide,topic):
    prompt = f"Generate slide titles for the topic '{topic}'. The number of slides should be between {MinSlide} and {MaxSlide}."
    msg = [
        {"role": "system", "content":prompt},
        {"role":"user", "content":topic}
    ]
    response = requests.post(
        url = "https://openrouter.ai/api/v1/chat/completions",
        headers = {"Authorization": f"Bearer {OPENROUTER_API_KEY}"},
        data = json.dumps({
            "messages": msg,
            "model": "openai/gpt-4o-mini-2024-07-18"
        })
    )
    resp = response.json()['choices'][0]['message']['content'].split("\n")
    return resp

def generate_slide_content(MaxLine,slide_title):
    prompt = f"Generate content for the slide: '{slide_title}'.And limit the content every page to {MaxLine} lines."
    msg = [
        {"role": "system", "content":prompt},
        {"role":"user", "content":slide_title}
    ]
    response = requests.post(
        url = "https://openrouter.ai/api/v1/chat/completions",
        headers = {"Authorization": f"Bearer {OPENROUTER_API_KEY}"},
        data = json.dumps({
            "messages": msg,
            "model": "openai/gpt-4o-mini-2024-07-18"
        })
    )
    resp = response.json()['choices'][0]['message']['content']
    return resp

def create_presentation(topic, slide_titles, slide_contents)->str:
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content
        #slide.shapes.placeholders[1].text = json.dumps(slide_content)

        # Customize font size for titles and content
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE
    # make sure no invalid characters in filename
    t = "".join([c for c in topic if c.isalpha() or c.isdigit() or c==' ']).rstrip()
    t = t.replace(" ", "_")
    prs.save(f"./{t}.pptx")
    return f"./{t}.pptx"

# to call the function, use the following code
# link = get_ppt_download_link(PATH_TO_PPT)
# st.markdown(link, unsafe_allow_html=True)  
# then you can have a hyperlink to download the file
def get_ppt_download_link(path, filename="presentation.pptx"):
    with open(path, "rb") as file:
        ppt_contents = file.read()
    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{path}",filename>Download the PowerPoint Presentation</a>'

# return url for button download
# to call the function, use the following code
# url = get_ppt_download_url(PATH_TO_PPT,PPT_NAME)
# st.link_button("Download Presentation", url)
def get_ppt_download_url(path):
    with open(path, "rb") as file:
        ppt_contents = file.read()
    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}'

# Convert PowerPoint file to PDF, need know the path 
def ppt_to_pdf(ppt_file_path, pdf_file_path):
    # Ensure the file paths are absolute
    ppt_file_path = os.path.abspath(ppt_file_path)
    pdf_file_path = os.path.abspath(pdf_file_path)

    # Check if the PPT file exists
    if not os.path.exists(ppt_file_path):
        raise FileNotFoundError(f"The file {ppt_file_path} does not exist.")

    # Create PowerPoint application object
    comtypes.CoInitialize()
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1

    try:
        # Open the presentation
        presentation = powerpoint.Presentations.Open(ppt_file_path)
        # Save as PDF
        presentation.SaveAs(pdf_file_path, 32)  # 32 is the formatType for PDF
        # Close the presentation
        presentation.Close()
    except Exception as e:
        print(f"An error occurred: {e}")
    finally:
        # Quit PowerPoint
        powerpoint.Quit()
        comtypes.CoUninitialize()

# use to display pdf file in web from path
def empdf(pdf_file_path):
    with open(pdf_file_path, "rb") as f:
        base64_pdf = base64.b64encode(f.read()).decode('utf-8')
    pdf_display = f'<iframe src="data:application/pdf;base64,{base64_pdf}" width="100%" height="800px" title="output_pdf"></iframe> '
    return pdf_display
# use to display pdf file in web from file
def embed_pdf(file):
    base64_pdf = base64.b64encode(file.read()).decode('utf-8')
    pdf_display = f'<iframe src="data:application/pdf;base64,{base64_pdf}" width="100%" height="800px"></iframe>'
    return pdf_display

def main():
    # web page title
    st.set_page_config(page_title="PowerPoint Presentation Generator",layout="wide")
    st.title("PowerPoint Presentation Generator")
    # Hide the menu
    st.markdown("""
    <style>
        .reportview-container {
            margin-top: -2em;
        }
        #MainMenu {visibility: hidden;}
        .stDeployButton {display:none;}
        header {visibility: hidden;}
        footer {visibility: hidden;}
        #stDecoration {display:none;}
    </style>
""", unsafe_allow_html=True)
    # Create two columns with different proportions
    col1, col2 = st.columns([3,2])
    txt = wd = df = None
    pf = md = ppt = None

    # defoult values
    if "MinSlide" not in st.session_state:
        st.session_state.MinSlide = 3
    if "MaxSlide" not in st.session_state:
        st.session_state.MaxSlide = 10
    if "MaxLine" not in st.session_state:
        st.session_state.MaxLine = 10
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    if "user_input" not in st.session_state:
        st.session_state.user_input = ""
    if "outputpath" not in st.session_state:
        st.session_state.outputpath = None
    if "generatedpath" not in st.session_state:
        st.session_state.generatedpath = None

    # popup window when generating presentation
    my_modal = Modal(title="Running Status", key="modal_key", max_width=600)

    with col2: # Right column
        topic = st.text_input("Enter the topic for your presentation:",help="mandatory")
        if topic == "":
            st.warning("Please enter a topic for your presentation.")
        # Create a row with three columns for the number inputs
        col_min, col_max, col_line = st.columns(3)
        with col_min:
            st.session_state.MinSlide = st.number_input("Min page num:", min_value=0, max_value=100, value=3, step=1)
        with col_max:
            st.session_state.MaxSlide = st.number_input("Max page num:", min_value=0, max_value=100, value=10, step=1)
        with col_line:
            st.session_state.MaxLine = st.number_input("MaxLine per page:", min_value=0, max_value=20, value=10, step=1)
        uploaded_file = st.file_uploader("Choose a file",type=['md','docx','csv','txt'])
        # New input area for user
        user_input = st.chat_input("If you have more requirements to state:")
        col_gen, col_exit = st.columns([1,1])
        with col_gen:
            GenerateButton = st.button("Generate Presentation")
        with col_exit:
            # the last ppt file in local path will be removed
            Exit = st.button("Exit Application And Clear Cache",type="primary")
        if st.session_state.generatedpath is not None:
            url = get_ppt_download_url(st.session_state.generatedpath)
            # st.link_button("Download Presentation", url) # not working
            # get the filename from the path after /
            filename = st.session_state.generatedpath.split("/")[-1]
            # style the download button
            st.markdown(f"""
            <a href="{url}" download={filename} style="
                display: inline-block;
                padding: 0.5em 1em;
                font-size: 1em;
                color: black;
                background-color: white;
                border: 1px solid green;
                border-radius: 0.25em;
                text-decoration: none;
                cursor: pointer;
            ">Download Presentation</a>
        """, unsafe_allow_html=True)
        if Exit:
            os.remove(st.session_state.generatedpath)  #clear cache file
            st.stop()
            # shut down the app

        # Add checkboxes for user to select options
        show_input_file = st.checkbox("Display Input File", value=True)
        show_output_preview = st.checkbox("Display Output Preview", value=True)
        st.markdown("##### Chat History")
        messages = st.container(height=300, key="chat_history_display")
        # if topic changes , we have to clear the chat history
        # first cell store current topic
        if st.session_state.chat_history != [] and st.session_state.chat_history[0] != "":
            for message in st.session_state.chat_history[1:]:
                messages.chat_message("user").write(message)
        # according to the file type, pass the content to variables
        if uploaded_file is not None:
            if uploaded_file.type == 'text/csv':
                df = pd.read_csv(uploaded_file)
            elif uploaded_file.type == 'application/octet-stream':
                md = uploaded_file.getvalue().decode("utf-8")
            elif uploaded_file.type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                wd = docx2txt.process(uploaded_file)
            elif uploaded_file.type == 'text/plain':
                txt = uploaded_file.getvalue().decode("utf-8")
            #** debug
            # elif uploaded_file.type == 'application/pdf':
            #     pf = uploaded_file
            # elif uploaded_file.type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            #     ppt = uploaded_file
            else:
                st.write("File not supported")
        
    with col1: # Left column
        if show_output_preview: # Display the output ppt preview
            if st.session_state.outputpath is not None:
                st.markdown("## Output Preview")
                ppt_to_pdf(st.session_state.outputpath, st.session_state.outputpath.replace(".pptx", ".pdf"))
                st.markdown(empdf(st.session_state.outputpath.replace(".pptx", ".pdf")), unsafe_allow_html=True)
                os.remove(st.session_state.outputpath.replace(".pptx", ".pdf"))
                st.session_state.generatedpath = st.session_state.outputpath
                st.session_state.outputpath = None

        if show_input_file: # Display the input file preview
            if (txt is not None) or (wd is not None) or (df is not None) or (md is not None):
                st.markdown("## Input Preview")
            if txt is not None:
                st.text_area("TXT File Content", txt, height=400)
            if wd is not None:
                st.write(wd)
            if df is not None:
                st.write(df)
            if md is not None:
                st.markdown(md)
            #** debug
            # if pf is not None:
            #     st.markdown(embed_pdf(pf), unsafe_allow_html=True)
            # if ppt is not None:
            #     # Save the uploaded file to a temporary location
            #     with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_ppt:
            #         tmp_ppt.write(ppt.getbuffer())
            #         tmp_ppt_path = tmp_ppt.name
            #         tmp_pdf_path = tmp_ppt_path.replace(".pptx", ".pdf")
            #     # Convert the PowerPoint file to pdf
            #     ppt_to_pdf(tmp_ppt_path, tmp_pdf_path)
            #     # Display the PDF file
            #     st.markdown(empdf(tmp_pdf_path), unsafe_allow_html=True)
            #     # Clean up the temporary PPTX file
            #     os.remove(tmp_ppt_path)
            #     os.remove(tmp_pdf_path) 
            

    # if GenerateButton and topic:
    if (GenerateButton or user_input) and topic:
        # every time user click the button, the local temp file ppt will be removed
        if st.session_state.generatedpath is not None:
                os.remove(st.session_state.generatedpath)
                st.session_state.generatedpath = None
        if user_input != None:
                # chat_history contains topic in first line
                if (st.session_state.chat_history == []) or topic != st.session_state.chat_history[0]:
                    st.session_state.chat_history = []
                    st.session_state.chat_history.append(topic)
                st.session_state.chat_history.append(f"You: {user_input}")
        # popup window to show the generating status
        with my_modal.container():
            with st.spinner("Generating presentation... Please wait."):
                    #** debug
                    # time.sleep(2)
    
                    # main part to generate the presentation
                    
                    # # you should call df, txt, wd, md to generate the presentation
                    # # chat_history is also need to be called in prompt
    
                    # # example
                    # slide_titles = generate_slide_titles(st.session_state.MinSlide,st.session_state.MaxSlide,topic)
                    # filtered_slide_titles= [item for item in slide_titles if item.strip() != '']
                    # #print("Slide Title: ", filtered_slide_titles)
    
                    # slide_contents = [generate_slide_content(st.session_state.MaxLine,title) for title in filtered_slide_titles]
                    # #print("Slide Contents: ", slide_contents)
                    # create_presentation(topic, filtered_slide_titles, slide_contents)
                    # #print("Presentation generated successfully!")
                    # st.markdown(get_ppt_download_link(topic), unsafe_allow_html=True)
                    
                    #** debug 
                    filtered_slide_titles = ['Introduction', 'Methodology', 'Results', 'Conclusion']
                    slide_contents = ['This is the introduction slide.', 'This is the methodology slide.', 'This is the results slide.', 'This is the conclusion slide.']
                    st.session_state.outputpath = create_presentation(topic, filtered_slide_titles, slide_contents)
                    st.session_state.generatedpath = st.session_state.outputpath
                    st.success("Presentation generated successfully!")
            if show_output_preview:
                    st.warning("Generating preview... Please wait.")
            my_modal.close() # Close the popup window automatically after generating the presentation
    
if __name__ == "__main__":
    main()


